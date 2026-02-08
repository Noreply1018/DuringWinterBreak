"""
提取Doc目录下的Word和PPT文件内容，转换为Markdown格式
"""

import os
from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

def extract_docx_to_markdown(docx_path: str) -> str:
    """
    提取Word文档内容并转换为Markdown格式
    """
    doc = Document(docx_path)
    md_content = []
    
    # 添加文件标题
    filename = Path(docx_path).stem
    md_content.append(f"# {filename}\n")
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        # 根据段落样式处理标题级别
        style_name = para.style.name if para.style else ""
        
        if "Heading 1" in style_name or "标题 1" in style_name:
            md_content.append(f"## {text}\n")
        elif "Heading 2" in style_name or "标题 2" in style_name:
            md_content.append(f"### {text}\n")
        elif "Heading 3" in style_name or "标题 3" in style_name:
            md_content.append(f"#### {text}\n")
        elif "Title" in style_name or "标题" in style_name:
            md_content.append(f"# {text}\n")
        else:
            md_content.append(f"{text}\n")
    
    # 处理表格
    for table in doc.tables:
        md_content.append("\n")
        
        # 获取表头
        if table.rows:
            header_row = table.rows[0]
            headers = [cell.text.strip().replace('\n', ' ') for cell in header_row.cells]
            md_content.append("| " + " | ".join(headers) + " |")
            md_content.append("| " + " | ".join(["---"] * len(headers)) + " |")
            
            # 处理数据行
            for row in table.rows[1:]:
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                md_content.append("| " + " | ".join(cells) + " |")
        
        md_content.append("\n")
    
    return "\n".join(md_content)


def extract_pptx_to_markdown(pptx_path: str) -> str:
    """
    提取PPT文档内容并转换为Markdown格式
    """
    prs = Presentation(pptx_path)
    md_content = []
    
    # 添加文件标题
    filename = Path(pptx_path).stem
    md_content.append(f"# {filename}\n")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        md_content.append(f"\n---\n\n## 幻灯片 {slide_num}\n")
        
        for shape in slide.shapes:
            # 处理文本框和形状中的文本
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # 检查是否为标题
                        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                            placeholder_type = shape.placeholder_format.type
                            # 1 = TITLE, 2 = BODY, 3 = CENTER_TITLE
                            if placeholder_type in [1, 3]:  # TITLE or CENTER_TITLE
                                md_content.append(f"### {text}\n")
                            else:
                                # 处理列表项
                                level = paragraph.level if paragraph.level else 0
                                indent = "  " * level
                                if paragraph.text.strip():
                                    md_content.append(f"{indent}- {text}\n")
                        else:
                            md_content.append(f"{text}\n")
            
            # 处理表格
            if shape.has_table:
                table = shape.table
                md_content.append("\n")
                
                if len(table.rows) > 0:
                    # 表头
                    header_cells = [cell.text.strip().replace('\n', ' ') for cell in table.rows[0].cells]
                    md_content.append("| " + " | ".join(header_cells) + " |")
                    md_content.append("| " + " | ".join(["---"] * len(header_cells)) + " |")
                    
                    # 数据行
                    for row in list(table.rows)[1:]:
                        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                        md_content.append("| " + " | ".join(cells) + " |")
                
                md_content.append("\n")
    
    return "\n".join(md_content)


def main():
    doc_dir = Path(r"d:\Desktop\寒假做五个项目\影像裁剪小工具\Doc")
    
    # 处理所有Word文件
    for docx_file in doc_dir.glob("*.docx"):
        print(f"正在处理 Word 文件: {docx_file.name}")
        try:
            md_content = extract_docx_to_markdown(str(docx_file))
            output_path = doc_dir / f"{docx_file.stem}.md"
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(md_content)
            print(f"  -> 已保存为: {output_path.name}")
        except Exception as e:
            print(f"  -> 处理失败: {e}")
    
    # 处理所有PPT文件
    for pptx_file in doc_dir.glob("*.pptx"):
        print(f"正在处理 PPT 文件: {pptx_file.name}")
        try:
            md_content = extract_pptx_to_markdown(str(pptx_file))
            output_path = doc_dir / f"{pptx_file.stem}.md"
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(md_content)
            print(f"  -> 已保存为: {output_path.name}")
        except Exception as e:
            print(f"  -> 处理失败: {e}")
    
    print("\n提取完成!")


if __name__ == "__main__":
    main()
